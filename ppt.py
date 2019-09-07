#coding:utf-8
import sys

reload(sys)
sys.setdefaultencoding("utf-8")
import requests
import json
import csv,time,os
import xlwt
from xlutils.copy import copy
from xlrd import open_workbook
import winsound
import pymysql
from win32com.client import Dispatch
from wxpy import *
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt #Inches
from pptx.enum.chart import XL_LEGEND_POSITION
# from docx.shared import RGBColor
from pptx.dml.color import RGBColor
import win32com
from pptx.dml.color import ColorFormat
from pptx.text.fonts import FontFiles


def zhuaqu():
    # 爬虫地址
    url = "http://ces.sino-life.com:7001/SL_CES/marketingChannel/queryMarketCollectList.do"
    headers = {
    "Host": "ces.sino-life.com:7001",
    "Connection": "keep-alive",
    "Content-Length": "14",
    "Accept": "application/json, text/javascript, */*; q=0.01",
    "X-Requested-With": "XMLHttpRequest",
    "SF_AJAX": "true",
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/76.0.3809.100 Safari/537.36",
    "Content-Type": "application/x-www-form-urlencoded",
    "Origin": "http://ces.sino-life.com:7001",
    "Accept-Encoding": "gzip, deflate",
    "Accept-Language": "zh-CN,zh;q=0.9",
    "Cookie":"""_sf_profile="63757272656e74556964:7765692e7a68656e673030364073696e6f2d6c6966652e636f6d@70726f66696c65496e6974466c6167:37@70726f66696c65436667466c6167:30"; JSESSIONID=F4B342ACFD9F88621B30384F40B508EE; SF_LOGIN_TIME="2019/09/04 14-24-20@2019/09/05 07-44-28"; BIGipServerCES_PRD_POOL_NEW=3129911488.9392.0000""",
    }
    data = {"channelType":"01","branchName":"葫芦岛中心支公司(862114)","startDate":"2019-09-5","endDate":"2019-09-5","productCode":"","page":"1","rows":"100"}
    response = requests.post(url=url, headers=headers,data=data, verify=False)
    jsonobj = json.loads(response.text)
    json_list =  jsonobj["rows"]

    return json_list

    # for jsonobj2 in json_list:
    #     name = jsonobj2['agentName']
    #     number = jsonobj2["agentNo"]
    #     money = jsonobj2["standardPrem"]
    #     place = jsonobj2["branchName"]
    #     xianzhong =jsonobj2["productName"]

def sql_insert():
    values = zhuaqu()
    dict_obj = {
        "I00665000001": "本级",
        "I00665000002": "本级",
        "I00665001395": "本级",
        "I00665002293": "本级",
        "I000000000230848": "本级",
        "I000000000394104": "本级",
        "I00665000919": "建昌",
        "I000000000254136": "建昌",
        "I000000000342993": "建昌",
        "I000000000389946": "建昌",
        "I000000000712626": "建昌",
        "I00665003107": "二区",
        "I000000000051679": "二区",
        "I000000000321700": "二区",
        "I000000000372278": "二区",
        "I000000000374813": "二区",
        "I00665000920": "一区",
        "I00665002691": "一区",
        "I00665002761": "一区",
        "I00665002832": "一区",
        "I00665003197": "一区",
        "I000000000025399": "一区",
        "I000000000301265": "一区",
        "I00665001015":"兴城",
        "I000000000268655":"兴城"

    }

    print(values)
    for date in values:
        db = db_obj()
        cursor = db.cursor()
        key = date["zoneCode"]
        if key in dict_obj:
            yingfu = dict_obj[key]
        sql = 'insert into fdjson.customer(`name`,`number`,`money`,`desc`,`product`,`code`) values (%s,%s,%s,%s,%s,%s)'
        if "富德生命康健无忧重大疾病保险" in date["productName"]:
            date["productName"] = "康健无忧"
        elif "富德生命鑫财富年金保险" in date["productName"]:
            date["productName"] = "鑫财富"
        elif"富德生命安行无忧两全保险" in date["productName"]:
            date["productName"] = "安行无忧"
        else:
            date["productName"] = ""
        cursor.execute(sql, (date["agentName"], date["agentNo"], date["standardPrem"], yingfu, date["productName"], date["zoneCode"]))
        db.commit()
        # work_sheet.write(row, 0, date["agentName"])
        # work_sheet.write(row, 1, date["agentNo"])
        # work_sheet.write(row, 2,date["standardPrem"])
        # work_sheet.write(row, 3,date["branchName"])
        # work_sheet.write(row, 4,date["productName"])
        # work_sheet.write(row, 5,date["zoneCode"])

        # row += 1
        # n+=1
    # print row
    # if row > rows_num:
    #     duration = 4000  # millisecond
    #     freq = 700  # Hz
    #     winsound.Beep(freq, duration)
    # excel.save("Python.xls")  # xlwt对象的保存方法，这时便覆盖掉了原来的excel
    #执行以下
    db = db_obj()
    sql = "select name,`desc`,sum(money),(case when product = '' or product is null then '康健无忧' else  product end) from customer GROUP BY `name`, `desc` ORDER BY MAX(id) "
    cursor = db.cursor()
    cursor.execute(sql)
    # 获取所有记录列表
    results = cursor.fetchall()

    rexcel = open_workbook("sujuku.xls",formatting_info=False)  # 用wlrd提供的方法读取一个excel文件
    excel = copy(rexcel)  # 用xlutils提供的copy方法将xlrd的对象转化为xlwt的对象
    work_sheet = excel.get_sheet(0)  # 用xlwt对象的方法获得要操作的sheet
    work_sheet.write(0, 0, u'姓名')
    work_sheet.write(0, 1, u'营服')
    work_sheet.write(0, 2, u'保费')
    work_sheet.write(0, 3, u'险种')

    row = 0
    rows_num = rexcel.sheets()[0].nrows
    print rows_num
    n = 0
    try:
        for x in results:
            row += 1
            # print x[0],x[1],x[2]


            work_sheet.write(row, 0, u'%s'%str(x[0]))
            work_sheet.write(row, 1,u'%s'%str(x[1]))
            work_sheet.write(row, 2,u'%s'%str(x[2]))
            work_sheet.write(row, 3, u'%s' % str(x[3]))

        ppt(x[0],x[1],x[3],x[2])

        excel.save("sujuku.xls")
        print "完成"
    except Exception as a:
        print a
def sql_dele():
    db = db_obj()
    sql ="TRUNCATE TABLE  customer ;"
    cursor = db.cursor()
    cursor.execute(sql)

def db_obj():
    # db = pymysql.connect("localhost", "root", "123456", "fdjson")
    conn = pymysql.connect(
        host="localhost",
        user='root',
        passwd='123456',
        charset='utf8',
        db="fdjson",
        use_unicode=False
    )
    return conn


def ppt(name,yingfu,xz,money):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)



    left = Inches(0)# 设置第一张图片的left，top
    top = Inches(0)
    height = Inches(7.5)
    weight = Inches(10)
    slide.shapes.add_picture("muban.jpg",left,top, height=height,width=weight)
    left = Inches(2.1)
    top = Inches(3.8)
    height = Inches(2)
    weight = Inches(4)
    textbox = slide.shapes.add_textbox(left, top,width=weight,height=height)  # left，top为相对位置，width，height为文本框大小
    tf = textbox.text_frame
    # tf.text = '祝贺建昌田玉坤伙伴\n      喜签康健无忧\n            6742元'  # 文本框中文字

    p = tf.add_paragraph()
    p.text = "祝贺%s%s伙伴\n      喜签%s\n            %s元"%(yingfu,name,xz,money)
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = RGBColor(250, 250, 0)
    prs.save("hebao.pptx")
    ppt_root = jpg_root = sys.path[0] + "\\"
    powerpoint = win32com.client.Dispatch('PowerPoint.Application')
    powerpoint.Visible = True
    ppt_path = ppt_root + "hebao.pptx"
    ppt = powerpoint.Presentations.Open(ppt_path)
    ppt.SaveAs(jpg_root + "hebao.pptx".rsplit('.')[0] + '.jpg', 17)
    ppt.Close()







if __name__ == '__main__':
    # shuju = zhuaqu()

    while 1:
        sql_insert()
        sql_dele()
        time.sleep(2)




